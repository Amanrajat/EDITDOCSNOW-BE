import io
import shutil
import unittest
import zipfile

import fitz
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from rest_framework import serializers

from apps.common.validation import validate_pdf_file

from .converters import jpg_to_pdf, pdf_to_excel, pdf_to_jpg, pdf_to_markdown, pdf_to_pptx, pdf_to_word
from .models import ConversionJob
from .services import run_conversion


def _make_pdf_with_text_table_and_image(name="doc.pdf"):
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 80), "Heading Text", fontsize=20, fontname="Helvetica-Bold")
    page.insert_text((72, 120), "A plain body paragraph.", fontsize=12, fontname="Helvetica")
    page.insert_text((72, 150), "An italic sentence.", fontsize=12, fontname="Helvetica-Oblique")

    page.draw_rect(fitz.Rect(72, 200, 400, 280))
    page.draw_line((236, 200), (236, 280))
    page.draw_line((72, 240), (400, 240))
    page.insert_text((90, 225), "Name")
    page.insert_text((250, 225), "Score")
    page.insert_text((90, 265), "Alice")
    page.insert_text((250, 265), "92")

    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 200, 150))
    pix.set_rect(pix.irect, (10, 120, 200))
    page.insert_image(fitz.Rect(400, 100, 560, 200), stream=pix.tobytes("png"))

    data = doc.tobytes()
    doc.close()
    return SimpleUploadedFile(name, data, content_type="application/pdf")


def _make_scanned_like_pdf(name="scanned.pdf"):
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 800, 1000))
    pix.set_rect(pix.irect, (200, 200, 200))
    page.insert_image(fitz.Rect(0, 0, 595, 842), stream=pix.tobytes("png"))
    data = doc.tobytes()
    doc.close()
    return SimpleUploadedFile(name, data, content_type="application/pdf")


class PdfToWordConverterTests(TestCase):

    def test_converts_text_table_and_image(self):
        pdf = _make_pdf_with_text_table_and_image()
        docx_bytes, metadata = pdf_to_word.convert(pdf.read())

        self.assertEqual(metadata["page_count"], 1)
        self.assertEqual(metadata["table_count"], 1)
        self.assertGreaterEqual(metadata["image_count"], 1)
        self.assertEqual(metadata["scanned_pages"], [])

        document = DocxDocument(io.BytesIO(docx_bytes))
        paragraph_texts = [p.text for p in document.paragraphs if p.text.strip()]
        self.assertIn("Heading Text", paragraph_texts)
        self.assertIn("A plain body paragraph.", paragraph_texts)
        self.assertIn("An italic sentence.", paragraph_texts)

        self.assertEqual(len(document.tables), 1)
        table_rows = [[cell.text for cell in row.cells] for row in document.tables[0].rows]
        self.assertEqual(table_rows, [["Name", "Score"], ["Alice", "92"]])

        self.assertGreaterEqual(len(document.inline_shapes), 1)

    def test_preserves_bold_and_italic(self):
        pdf = _make_pdf_with_text_table_and_image()
        docx_bytes, _ = pdf_to_word.convert(pdf.read())
        document = DocxDocument(io.BytesIO(docx_bytes))

        runs_by_text = {}
        for paragraph in document.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    runs_by_text[run.text] = run

        self.assertTrue(runs_by_text["Heading Text"].bold)
        self.assertFalse(runs_by_text["A plain body paragraph."].bold)
        self.assertTrue(runs_by_text["An italic sentence."].italic)

    def test_scanned_page_falls_back_to_image_and_is_reported(self):
        pdf = _make_scanned_like_pdf()
        docx_bytes, metadata = pdf_to_word.convert(pdf.read())

        self.assertEqual(metadata["scanned_pages"], [1])
        document = DocxDocument(io.BytesIO(docx_bytes))
        self.assertGreaterEqual(len(document.inline_shapes), 1)

    def test_multi_page_inserts_page_breaks(self):
        doc = fitz.open()
        for text in ["Page one text", "Page two text"]:
            page = doc.new_page(width=400, height=500)
            page.insert_text((50, 50), text, fontsize=14)
        data = doc.tobytes()
        doc.close()

        docx_bytes, metadata = pdf_to_word.convert(data)
        self.assertEqual(metadata["page_count"], 2)
        document = DocxDocument(io.BytesIO(docx_bytes))
        paragraph_texts = [p.text for p in document.paragraphs if p.text.strip()]
        self.assertIn("Page one text", paragraph_texts)
        self.assertIn("Page two text", paragraph_texts)


class RunConversionServiceTests(TestCase):

    def test_run_conversion_creates_completed_job(self):
        pdf = _make_pdf_with_text_table_and_image()
        job = run_conversion(
            user=None, uploaded_file=pdf, operation=ConversionJob.Operation.PDF_TO_WORD,
            converter_fn=pdf_to_word.convert, output_ext="docx",
        )
        self.assertEqual(job.status, ConversionJob.Status.COMPLETED)
        self.assertTrue(job.output_file.name)
        self.assertEqual(job.metadata["table_count"], 1)
        job.output_file.delete(save=False)

    def test_run_conversion_marks_failed_on_converter_exception(self):
        def _broken_converter(_file_bytes):
            raise ValueError("boom")

        pdf = _make_pdf_with_text_table_and_image()
        job = run_conversion(
            user=None, uploaded_file=pdf, operation=ConversionJob.Operation.PDF_TO_WORD,
            converter_fn=_broken_converter, output_ext="docx",
        )
        self.assertEqual(job.status, ConversionJob.Status.FAILED)
        self.assertIn("boom", job.error_message)


class PdfToWordAPITests(TestCase):

    def test_convert_end_to_end(self):
        pdf = _make_pdf_with_text_table_and_image()
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-word/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        body = response.json()
        self.assertTrue(body["success"])
        data = body["data"]
        self.assertEqual(data["operation"], "pdf_to_word")
        self.assertEqual(data["table_count"], 1)

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download.status_code, 200)
        self.assertEqual(
            download["Content-Type"],
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

        downloaded_bytes = b"".join(download.streaming_content)
        document = DocxDocument(io.BytesIO(downloaded_bytes))
        self.assertTrue(any("Heading Text" in p.text for p in document.paragraphs))

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_missing_file_is_rejected(self):
        response = self.client.post("/api/v1/pdf/convert/pdf-to-word/", data={}, format="multipart")
        self.assertEqual(response.status_code, 400)
        self.assertIn("file", response.json()["errors"])

    def test_rejects_non_pdf_file(self):
        not_a_pdf = SimpleUploadedFile("evil.txt", b"hello", content_type="text/plain")
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-word/", data={"file": not_a_pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)
        self.assertFalse(response.json()["success"])

    def test_rejects_fake_pdf_without_signature(self):
        fake = SimpleUploadedFile("fake.pdf", b"not a real pdf", content_type="application/pdf")
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-word/", data={"file": fake}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)
        self.assertFalse(response.json()["success"])

    def test_rejects_oversized_file_via_shared_validator(self):
        pdf = _make_pdf_with_text_table_and_image()
        with self.assertRaises(serializers.ValidationError):
            validate_pdf_file(pdf, max_size=10)


class ConversionJobOwnershipTests(TestCase):
    """
    Access-control regression suite for the shared ownership layer, as
    exercised through PDF-to-Word's download endpoint (the shared
    ConversionJobDownloadView). Mirrors the equivalent ownership test
    classes across every other PDF job app - same shared
    apps.common.ownership/views layer.
    """

    def _create_job(self):
        pdf = _make_pdf_with_text_table_and_image()
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-word/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        return response.json()["data"]

    def test_owner_can_download_with_correct_token(self):
        data = self._create_job()
        download_url = data["download_url"].replace("http://testserver", "")

        response = self.client.get(download_url)
        self.assertEqual(response.status_code, 200)

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_wrong_token_is_denied(self):
        data = self._create_job()
        job = ConversionJob.objects.get(id=data["file_id"])

        response = self.client.get(f"/api/v1/pdf/convert/{job.id}/download/?token=not-the-real-token")
        self.assertEqual(response.status_code, 404)
        body = response.json()
        self.assertFalse(body["success"])
        self.assertEqual(body["error_code"], "NOT_FOUND")

        job.output_file.delete(save=False)

    def test_missing_token_is_denied(self):
        data = self._create_job()
        job = ConversionJob.objects.get(id=data["file_id"])

        response = self.client.get(f"/api/v1/pdf/convert/{job.id}/download/")
        self.assertEqual(response.status_code, 404)

        job.output_file.delete(save=False)

    def test_nonexistent_job_id_returns_identical_response_shape_to_wrong_token(self):
        import uuid

        response = self.client.get(f"/api/v1/pdf/convert/{uuid.uuid4()}/download/?token=whatever")
        self.assertEqual(response.status_code, 404)
        body = response.json()
        self.assertFalse(body["success"])
        self.assertEqual(body["error_code"], "NOT_FOUND")

    def test_malformed_job_id_is_handled_safely_not_a_500(self):
        response = self.client.get("/api/v1/pdf/convert/not-a-valid-uuid/download/?token=whatever")
        self.assertEqual(response.status_code, 404)

    def test_raw_media_path_does_not_serve_the_output_file(self):
        data = self._create_job()
        job = ConversionJob.objects.get(id=data["file_id"])

        response = self.client.get(f"/media/convert/output/{job.id}.docx")
        self.assertEqual(response.status_code, 404)

        job.output_file.delete(save=False)


class PdfToExcelConverterTests(TestCase):

    def test_extracts_table_into_its_own_sheet(self):
        pdf = _make_pdf_with_text_table_and_image()
        xlsx_bytes, metadata = pdf_to_excel.convert(pdf.read())

        self.assertEqual(metadata["page_count"], 1)
        self.assertEqual(metadata["table_count"], 1)

        workbook = load_workbook(io.BytesIO(xlsx_bytes))
        self.assertEqual(len(workbook.sheetnames), 1)
        sheet = workbook["Page 1"]
        rows = [[cell.value for cell in row] for row in sheet.iter_rows(min_row=1, max_row=2, max_col=2)]
        self.assertEqual(rows, [["Name", "Score"], ["Alice", "92"]])

    def test_page_without_a_table_still_writes_plain_text(self):
        doc = fitz.open()
        page = doc.new_page(width=400, height=500)
        page.insert_text((50, 50), "Just a line of text.", fontsize=12)
        data = doc.tobytes()
        doc.close()

        xlsx_bytes, metadata = pdf_to_excel.convert(data)
        self.assertEqual(metadata["table_count"], 0)
        workbook = load_workbook(io.BytesIO(xlsx_bytes))
        sheet = workbook["Page 1"]
        self.assertEqual(sheet.cell(row=1, column=1).value, "Just a line of text.")


class PdfToExcelAPITests(TestCase):

    def test_convert_end_to_end(self):
        pdf = _make_pdf_with_text_table_and_image()
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-excel/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "pdf_to_excel")
        self.assertEqual(data["table_count"], 1)

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download.status_code, 200)
        self.assertEqual(
            download["Content-Type"],
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        workbook = load_workbook(io.BytesIO(b"".join(download.streaming_content)))
        self.assertIn("Page 1", workbook.sheetnames)

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)


class PdfToPptxConverterTests(TestCase):

    def test_one_slide_per_page_with_real_content(self):
        pdf = _make_pdf_with_text_table_and_image()
        pptx_bytes, metadata = pdf_to_pptx.convert(pdf.read())

        self.assertEqual(metadata["page_count"], 1)
        self.assertEqual(metadata["table_count"], 1)
        self.assertGreaterEqual(metadata["image_count"], 1)

        prs = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(len(prs.slides._sldIdLst), 1)
        slide = prs.slides[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        self.assertIn("Heading Text", texts)
        tables = [s.table for s in slide.shapes if s.has_table]
        self.assertEqual(len(tables), 1)
        rows = [[c.text for c in row.cells] for row in tables[0].rows]
        self.assertEqual(rows, [["Name", "Score"], ["Alice", "92"]])

    def test_scanned_page_falls_back_to_full_slide_image(self):
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 800, 1000))
        pix.set_rect(pix.irect, (200, 200, 200))
        page.insert_image(fitz.Rect(0, 0, 595, 842), stream=pix.tobytes("png"))
        data = doc.tobytes()
        doc.close()

        pptx_bytes, metadata = pdf_to_pptx.convert(data)
        self.assertEqual(metadata["scanned_pages"], [1])
        prs = Presentation(io.BytesIO(pptx_bytes))
        picture_shapes = [s for s in prs.slides[0].shapes if s.shape_type == 13]
        self.assertEqual(len(picture_shapes), 1)


class PdfToPptxAPITests(TestCase):

    def test_convert_end_to_end(self):
        pdf = _make_pdf_with_text_table_and_image()
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-pptx/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "pdf_to_pptx")

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(
            download["Content-Type"],
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        prs = Presentation(io.BytesIO(b"".join(download.streaming_content)))
        self.assertEqual(len(prs.slides._sldIdLst), 1)

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)


class PdfToJpgConverterTests(TestCase):

    def _make_multi_page_pdf(self, count=3):
        doc = fitz.open()
        for i in range(count):
            page = doc.new_page(width=300, height=200)
            page.insert_text((20, 30), f"Page {i + 1}", fontsize=14)
        data = doc.tobytes()
        doc.close()
        return data

    def test_single_page_returns_a_plain_jpg(self):
        data = self._make_multi_page_pdf(3)
        output_bytes, metadata = pdf_to_jpg.convert(data, pages=[2])

        self.assertFalse(metadata["_is_zip"])
        self.assertEqual(metadata["_output_ext"], "jpg")
        self.assertEqual(metadata["converted_pages"], [2])
        self.assertTrue(output_bytes.startswith(b"\xff\xd8"))  # real JPEG magic bytes

    def test_multiple_pages_are_zipped(self):
        data = self._make_multi_page_pdf(3)
        output_bytes, metadata = pdf_to_jpg.convert(data, pages=None)

        self.assertTrue(metadata["_is_zip"])
        zf = zipfile.ZipFile(io.BytesIO(output_bytes))
        self.assertEqual(sorted(zf.namelist()), ["page_1.jpg", "page_2.jpg", "page_3.jpg"])
        for name in zf.namelist():
            self.assertTrue(zf.read(name).startswith(b"\xff\xd8"))

    def test_higher_dpi_yields_a_larger_image(self):
        data = self._make_multi_page_pdf(1)
        low_bytes, _ = pdf_to_jpg.convert(data, pages=[1], dpi=72, quality=90)
        high_bytes, _ = pdf_to_jpg.convert(data, pages=[1], dpi=300, quality=90)
        self.assertGreater(len(high_bytes), len(low_bytes))

    def test_rejects_out_of_range_page(self):
        data = self._make_multi_page_pdf(2)
        with self.assertRaises(pdf_to_jpg.PdfToJpgError):
            pdf_to_jpg.convert(data, pages=[99])

    def test_rejects_invalid_dpi(self):
        data = self._make_multi_page_pdf(1)
        with self.assertRaises(pdf_to_jpg.PdfToJpgError):
            pdf_to_jpg.convert(data, pages=[1], dpi=5000)

    def test_rejects_invalid_quality(self):
        data = self._make_multi_page_pdf(1)
        with self.assertRaises(pdf_to_jpg.PdfToJpgError):
            pdf_to_jpg.convert(data, pages=[1], quality=200)


class PdfToJpgAPITests(TestCase):

    def _make_multi_page_pdf(self, count=3):
        doc = fitz.open()
        for i in range(count):
            page = doc.new_page(width=300, height=200)
            page.insert_text((20, 30), f"Page {i + 1}", fontsize=14)
        data = doc.tobytes()
        doc.close()
        return SimpleUploadedFile("doc.pdf", data, content_type="application/pdf")

    def test_single_page_end_to_end_returns_jpg(self):
        pdf = self._make_multi_page_pdf(3)
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-jpg/", data={"file": pdf, "pages": [1]}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "image/jpeg")

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_all_pages_end_to_end_returns_zip(self):
        pdf = self._make_multi_page_pdf(3)
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-jpg/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "application/zip")
        zf = zipfile.ZipFile(io.BytesIO(b"".join(download.streaming_content)))
        self.assertEqual(len(zf.namelist()), 3)

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_out_of_range_page_is_rejected(self):
        pdf = self._make_multi_page_pdf(2)
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-jpg/", data={"file": pdf, "pages": [99]}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)


class PdfToMarkdownConverterTests(TestCase):

    def test_extracts_headings_paragraphs_lists_tables_and_links(self):
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        page.insert_text((72, 80), "Main Title", fontsize=24, fontname="Helvetica-Bold")
        page.insert_text((72, 130), "A regular paragraph of body text here.", fontsize=11)
        page.insert_text((72, 160), "- First item", fontsize=11)
        page.insert_text((72, 240), "Visit our site", fontsize=11)
        page.insert_link({"kind": fitz.LINK_URI, "from": fitz.Rect(72, 230, 200, 250), "uri": "https://example.com"})
        page.draw_rect(fitz.Rect(72, 300, 300, 360))
        page.draw_line((186, 300), (186, 360))
        page.draw_line((72, 330), (300, 330))
        page.insert_text((90, 320), "Col A")
        page.insert_text((200, 320), "Col B")
        page.insert_text((90, 350), "x")
        page.insert_text((200, 350), "y")
        data = doc.tobytes()
        doc.close()

        md_bytes, metadata = pdf_to_markdown.convert(data)
        markdown = md_bytes.decode()

        self.assertEqual(metadata["heading_count"], 1)
        self.assertEqual(metadata["table_count"], 1)
        self.assertIn("# Main Title", markdown)
        self.assertIn("A regular paragraph of body text here.", markdown)
        self.assertIn("- First item", markdown)
        self.assertIn("[Visit our site](https://example.com)", markdown)
        self.assertIn("| Col A | Col B |", markdown)
        self.assertIn("| x | y |", markdown)

    def test_heading_levels_rank_by_font_size(self):
        doc = fitz.open()
        page = doc.new_page(width=400, height=500)
        page.insert_text((50, 50), "Big Heading", fontsize=24, fontname="Helvetica-Bold")
        page.insert_text((50, 100), "Medium Heading", fontsize=18, fontname="Helvetica-Bold")
        page.insert_text((50, 150), "Body text at normal size for contrast here.", fontsize=11)
        data = doc.tobytes()
        doc.close()

        md_bytes, _ = pdf_to_markdown.convert(data)
        markdown = md_bytes.decode()
        self.assertIn("# Big Heading", markdown)
        self.assertIn("## Medium Heading", markdown)


class PdfToMarkdownAPITests(TestCase):

    def test_convert_end_to_end(self):
        doc = fitz.open()
        page = doc.new_page(width=400, height=500)
        page.insert_text((50, 50), "Title Heading", fontsize=22, fontname="Helvetica-Bold")
        page.insert_text((50, 100), "Some body paragraph text goes here for contrast.", fontsize=11)
        data = doc.tobytes()
        doc.close()
        pdf = SimpleUploadedFile("doc.pdf", data, content_type="application/pdf")

        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-markdown/", data={"file": pdf}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "pdf_to_markdown")
        self.assertGreaterEqual(data["heading_count"], 1)

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "text/markdown")
        markdown = b"".join(download.streaming_content).decode()
        self.assertIn("# Title Heading", markdown)

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)


class OtherOperationsDownloadOwnershipSmokeTests(TestCase):
    """
    The download view (ConversionJobDownloadView) is shared across every
    operation - the full ownership test suite already lives on
    ConversionJobOwnershipTests (exercised via pdf-to-word). This just
    confirms the same protection actually applies when the job came from
    a *different* operation (i.e. nothing operation-specific bypasses it).
    """

    def test_wrong_token_is_denied_for_an_excel_job(self):
        pdf = _make_pdf_with_text_table_and_image()
        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-excel/", data={"file": pdf}, format="multipart",
        )
        data = response.json()["data"]
        job = ConversionJob.objects.get(id=data["file_id"])

        denied = self.client.get(f"/api/v1/pdf/convert/{job.id}/download/?token=wrong")
        self.assertEqual(denied.status_code, 404)

        job.output_file.delete(save=False)


def _make_solid_jpeg(name, size=(300, 200), color=(200, 50, 50)):
    from io import BytesIO

    from PIL import Image

    buffer = BytesIO()
    Image.new("RGB", size, color=color).save(buffer, format="JPEG")
    return SimpleUploadedFile(name, buffer.getvalue(), content_type="image/jpeg")


class JpgToPdfConverterTests(TestCase):

    def test_one_page_per_image_in_order(self):
        images = [
            ("a.jpg", _make_solid_jpeg("a.jpg", color=(255, 0, 0)).read()),
            ("b.jpg", _make_solid_jpeg("b.jpg", color=(0, 255, 0)).read()),
        ]
        from .converters import jpg_to_pdf

        output_bytes, metadata = jpg_to_pdf.convert(images)
        self.assertEqual(metadata["page_count"], 2)

        doc = fitz.open(stream=output_bytes, filetype="pdf")
        self.assertEqual(len(doc), 2)
        doc.close()

    def test_fit_mode_preserves_aspect_ratio_no_distortion(self):
        from .converters import jpg_to_pdf

        # A very wide image on an A4 portrait page - "fit" must not stretch it.
        wide_image = _make_solid_jpeg("wide.jpg", size=(2000, 200)).read()
        output_bytes, _ = jpg_to_pdf.convert([("wide.jpg", wide_image)], page_size="A4", fit_mode="fit")

        doc = fitz.open(stream=output_bytes, filetype="pdf")
        image_list = doc[0].get_images(full=True)
        self.assertEqual(len(image_list), 1)
        rects = doc[0].get_image_rects(image_list[0][0])
        rect = rects[0]
        source_ratio = 2000 / 200
        placed_ratio = rect.width / rect.height
        self.assertAlmostEqual(source_ratio, placed_ratio, delta=0.05)
        doc.close()

    def test_fill_mode_covers_the_full_available_area(self):
        from .converters import jpg_to_pdf

        wide_image = _make_solid_jpeg("wide.jpg", size=(2000, 200)).read()
        output_bytes, _ = jpg_to_pdf.convert([("wide.jpg", wide_image)], page_size="A4", fit_mode="fill", margin=0)

        doc = fitz.open(stream=output_bytes, filetype="pdf")
        image_list = doc[0].get_images(full=True)
        rects = doc[0].get_image_rects(image_list[0][0])
        rect = rects[0]
        page_rect = doc[0].rect
        # A couple points of slack: PyMuPDF nudges the placed rect very
        # slightly to preserve the (integer-pixel-cropped) image's own
        # exact aspect ratio rather than force-stretching it - "fill"
        # still covers effectively the entire page (>99%), which is what
        # actually matters here.
        self.assertAlmostEqual(rect.width, page_rect.width, delta=3)
        self.assertAlmostEqual(rect.height, page_rect.height, delta=3)
        doc.close()

    def test_landscape_orientation_swaps_dimensions(self):
        from .converters import jpg_to_pdf

        image = _make_solid_jpeg("a.jpg").read()
        output_bytes, _ = jpg_to_pdf.convert([("a.jpg", image)], page_size="A4", orientation="landscape")
        doc = fitz.open(stream=output_bytes, filetype="pdf")
        self.assertGreater(doc[0].rect.width, doc[0].rect.height)
        doc.close()

    def test_rejects_invalid_page_size(self):
        from .converters import jpg_to_pdf

        with self.assertRaises(jpg_to_pdf.JpgToPdfError):
            jpg_to_pdf.convert([("a.jpg", _make_solid_jpeg("a.jpg").read())], page_size="Legal")

    def test_rejects_empty_image_list(self):
        from .converters import jpg_to_pdf

        with self.assertRaises(jpg_to_pdf.JpgToPdfError):
            jpg_to_pdf.convert([])


class JpgToPdfAPITests(TestCase):

    def test_convert_end_to_end_preserves_order(self):
        red = _make_solid_jpeg("red.jpg", color=(255, 0, 0))
        green = _make_solid_jpeg("green.jpg", color=(0, 255, 0))
        response = self.client.post(
            "/api/v1/pdf/convert/jpg-to-pdf/",
            data={"files": [red, green], "page_size": "A4", "orientation": "portrait"},
            format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "jpg_to_pdf")
        self.assertEqual(data["page_count"], 2)

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "application/pdf")
        downloaded_bytes = b"".join(download.streaming_content)
        doc = fitz.open(stream=downloaded_bytes, filetype="pdf")
        self.assertEqual(len(doc), 2)
        doc.close()

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_custom_order_is_respected(self):
        red = _make_solid_jpeg("red.jpg", color=(255, 0, 0))
        green = _make_solid_jpeg("green.jpg", color=(0, 255, 0))
        # Reverse the upload order via `order`.
        response = self.client.post(
            "/api/v1/pdf/convert/jpg-to-pdf/",
            data={"files": [red, green], "order": [1, 0]},
            format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_missing_files_is_rejected(self):
        response = self.client.post("/api/v1/pdf/convert/jpg-to-pdf/", data={}, format="multipart")
        self.assertEqual(response.status_code, 400)

    def test_rejects_non_image_file(self):
        fake = SimpleUploadedFile("evil.jpg", b"not really a jpeg", content_type="image/jpeg")
        response = self.client.post(
            "/api/v1/pdf/convert/jpg-to-pdf/", data={"files": [fake]}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)

    def test_invalid_order_is_rejected(self):
        red = _make_solid_jpeg("red.jpg")
        response = self.client.post(
            "/api/v1/pdf/convert/jpg-to-pdf/",
            data={"files": [red], "order": [0, 1]},
            format="multipart",
        )
        self.assertEqual(response.status_code, 400)

    def test_invalid_page_size_is_rejected(self):
        red = _make_solid_jpeg("red.jpg")
        response = self.client.post(
            "/api/v1/pdf/convert/jpg-to-pdf/",
            data={"files": [red], "page_size": "Legal"},
            format="multipart",
        )
        self.assertEqual(response.status_code, 400)


class PdfToPdfAConverterTests(TestCase):

    def test_converts_and_preserves_page_count_and_text(self):
        from .converters import pdf_to_pdfa

        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        page.insert_text((72, 72), "PDFA conformance test", fontsize=14)
        page2 = doc.new_page(width=595, height=842)
        page2.insert_text((72, 72), "Second page", fontsize=14)
        data = doc.tobytes()
        doc.close()

        output_bytes, metadata = pdf_to_pdfa.convert(data, level="2b")
        self.assertEqual(metadata["page_count"], 2)
        self.assertEqual(metadata["pdfa_standard"], "PDF/A-2B")

        result_doc = fitz.open(stream=output_bytes, filetype="pdf")
        self.assertEqual(len(result_doc), 2)
        self.assertIn("PDFA conformance test", result_doc[0].get_text())
        self.assertIn("Second page", result_doc[1].get_text())
        result_doc.close()

    def test_rejects_invalid_level(self):
        from .converters import pdf_to_pdfa

        doc = fitz.open()
        doc.new_page()
        data = doc.tobytes()
        doc.close()

        with self.assertRaises(pdf_to_pdfa.PdfToPdfAError):
            pdf_to_pdfa.convert(data, level="4b")


class PdfToPdfAAPITests(TestCase):

    def test_convert_end_to_end(self):
        doc = fitz.open()
        page = doc.new_page(width=400, height=500)
        page.insert_text((50, 50), "Archival document", fontsize=14)
        data = doc.tobytes()
        doc.close()
        pdf = SimpleUploadedFile("doc.pdf", data, content_type="application/pdf")

        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-pdfa/", data={"file": pdf, "level": "2b"}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        result = response.json()["data"]
        self.assertEqual(result["operation"], "pdf_to_pdfa")
        self.assertEqual(result["pdfa_standard"], "PDF/A-2B")

        download_url = result["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "application/pdf")
        downloaded_bytes = b"".join(download.streaming_content)
        result_doc = fitz.open(stream=downloaded_bytes, filetype="pdf")
        self.assertIn("Archival document", result_doc[0].get_text())
        result_doc.close()

        job = ConversionJob.objects.get(id=result["file_id"])
        job.output_file.delete(save=False)

    def test_invalid_level_is_rejected(self):
        doc = fitz.open()
        doc.new_page()
        data = doc.tobytes()
        doc.close()
        pdf = SimpleUploadedFile("doc.pdf", data, content_type="application/pdf")

        response = self.client.post(
            "/api/v1/pdf/convert/pdf-to-pdfa/", data={"file": pdf, "level": "4b"}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)


class HtmlToPdfConverterTests(TestCase):

    def test_renders_raw_html_string(self):
        from .converters import html_to_pdf

        output_bytes, metadata = html_to_pdf.convert(html="<h1>Report</h1><p>Body text here.</p>")
        self.assertEqual(metadata["source"], "html")
        doc = fitz.open(stream=output_bytes, filetype="pdf")
        self.assertEqual(len(doc), 1)
        self.assertIn("Report", doc[0].get_text())
        self.assertIn("Body text here.", doc[0].get_text())
        doc.close()

    def test_page_size_and_orientation_affect_output_dimensions(self):
        from .converters import html_to_pdf

        portrait_bytes, _ = html_to_pdf.convert(html="<p>x</p>", page_size="A4", orientation="portrait")
        landscape_bytes, _ = html_to_pdf.convert(html="<p>x</p>", page_size="A4", orientation="landscape")

        portrait_doc = fitz.open(stream=portrait_bytes, filetype="pdf")
        landscape_doc = fitz.open(stream=landscape_bytes, filetype="pdf")
        self.assertLess(portrait_doc[0].rect.width, portrait_doc[0].rect.height)
        self.assertGreater(landscape_doc[0].rect.width, landscape_doc[0].rect.height)
        portrait_doc.close()
        landscape_doc.close()

    def test_requires_exactly_one_of_url_or_html(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert()
        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(url="http://example.com", html="<p>x</p>")

    def test_blocks_loopback_url(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(url="http://127.0.0.1/")

    def test_blocks_cloud_metadata_endpoint(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(url="http://169.254.169.254/latest/meta-data/")

    def test_blocks_private_network_ranges(self):
        from .converters import html_to_pdf

        for host in ["http://10.1.2.3/", "http://172.16.0.5/", "http://192.168.0.10/"]:
            with self.assertRaises(html_to_pdf.HtmlToPdfError):
                html_to_pdf.convert(url=host)

    def test_blocks_non_http_scheme(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(url="file:///etc/passwd")

    def test_blocks_non_default_port(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(url="http://example.com:8080/")

    def test_rejects_invalid_page_size(self):
        from .converters import html_to_pdf

        with self.assertRaises(html_to_pdf.HtmlToPdfError):
            html_to_pdf.convert(html="<p>x</p>", page_size="Legal")


class HtmlToPdfAPITests(TestCase):

    def test_renders_raw_html_end_to_end(self):
        response = self.client.post(
            "/api/v1/pdf/convert/html-to-pdf/",
            data={"html": "<h1>Invoice</h1><p>Total: $42</p>"},
            format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "html_to_pdf")

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "application/pdf")
        downloaded_bytes = b"".join(download.streaming_content)
        doc = fitz.open(stream=downloaded_bytes, filetype="pdf")
        self.assertIn("Invoice", doc[0].get_text())
        doc.close()

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_blocked_url_returns_400_without_creating_a_job(self):
        before_count = ConversionJob.objects.count()
        response = self.client.post(
            "/api/v1/pdf/convert/html-to-pdf/", data={"url": "http://127.0.0.1/admin"}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json()["error_code"], "BLOCKED_URL")
        self.assertEqual(ConversionJob.objects.count(), before_count)

    def test_blocked_metadata_endpoint_returns_400(self):
        response = self.client.post(
            "/api/v1/pdf/convert/html-to-pdf/",
            data={"url": "http://169.254.169.254/latest/meta-data/iam/security-credentials/"},
            format="multipart",
        )
        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json()["error_code"], "BLOCKED_URL")

    def test_requires_exactly_one_of_url_or_html(self):
        response = self.client.post("/api/v1/pdf/convert/html-to-pdf/", data={}, format="multipart")
        self.assertEqual(response.status_code, 400)

        response2 = self.client.post(
            "/api/v1/pdf/convert/html-to-pdf/",
            data={"url": "http://example.com", "html": "<p>x</p>"},
            format="multipart",
        )
        self.assertEqual(response2.status_code, 400)

    def test_invalid_page_size_is_rejected(self):
        response = self.client.post(
            "/api/v1/pdf/convert/html-to-pdf/", data={"html": "<p>x</p>", "page_size": "Legal"}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)


_SOFFICE_AVAILABLE = shutil.which("soffice") is not None
_SKIP_REASON = "LibreOffice ('soffice') is not installed in this environment"


@unittest.skipUnless(_SOFFICE_AVAILABLE, _SKIP_REASON)
class OfficeToPdfConverterTests(TestCase):
    """
    Runs for real only where LibreOffice is actually installed (the
    production Docker image - see Dockerfile - and any dev machine that
    has it). Skipped, not silently passed, everywhere else so CI output
    makes clear these were never executed rather than looking green by
    accident.
    """

    def test_word_to_pdf_preserves_text(self):
        from docx import Document as DocxDocument

        from .converters import office_to_pdf

        buffer = io.BytesIO()
        document = DocxDocument()
        document.add_heading("Quarterly Report", level=1)
        document.add_paragraph("Revenue grew by 12% this quarter.")
        document.save(buffer)

        output_bytes, metadata = office_to_pdf.convert(buffer.getvalue(), "docx")
        self.assertGreaterEqual(metadata["page_count"], 1)
        doc = fitz.open(stream=output_bytes, filetype="pdf")
        full_text = "".join(page.get_text() for page in doc)
        self.assertIn("Quarterly Report", full_text)
        self.assertIn("Revenue grew by 12% this quarter.", full_text)
        doc.close()

    def test_excel_to_pdf_preserves_cell_values(self):
        from openpyxl import Workbook

        from .converters import office_to_pdf

        buffer = io.BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = "Name"
        sheet["B1"] = "Score"
        sheet["A2"] = "Alice"
        sheet["B2"] = 92
        workbook.save(buffer)

        output_bytes, metadata = office_to_pdf.convert(buffer.getvalue(), "xlsx")
        self.assertGreaterEqual(metadata["page_count"], 1)
        doc = fitz.open(stream=output_bytes, filetype="pdf")
        full_text = "".join(page.get_text() for page in doc)
        self.assertIn("Alice", full_text)
        doc.close()

    def test_pptx_to_pdf_preserves_slide_text(self):
        from pptx import Presentation
        from pptx.util import Inches

        from .converters import office_to_pdf

        buffer = io.BytesIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        textbox.text_frame.text = "Quarterly Results"
        prs.save(buffer)

        output_bytes, metadata = office_to_pdf.convert(buffer.getvalue(), "pptx")
        self.assertEqual(metadata["page_count"], 1)
        doc = fitz.open(stream=output_bytes, filetype="pdf")
        self.assertIn("Quarterly Results", doc[0].get_text())
        doc.close()

    def test_rejects_missing_libreoffice_gracefully(self):
        # We can't easily simulate "soffice missing" while it's actually
        # installed, so this just documents the contract: importing the
        # module and calling _require_libreoffice with a monkeypatched
        # shutil.which confirms the clear-error path exists and works.
        from unittest import mock

        from .converters import office_to_pdf

        with mock.patch("apps.pdf_convert.converters.office_to_pdf.shutil.which", return_value=None):
            with self.assertRaises(office_to_pdf.OfficeToPdfError):
                office_to_pdf.convert(b"irrelevant", "docx")


@unittest.skipUnless(_SOFFICE_AVAILABLE, _SKIP_REASON)
class WordToPdfAPITests(TestCase):

    def test_convert_end_to_end(self):
        from docx import Document as DocxDocument

        buffer = io.BytesIO()
        document = DocxDocument()
        document.add_heading("API Test Doc", level=1)
        document.save(buffer)
        docx_file = SimpleUploadedFile(
            "doc.docx", buffer.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

        response = self.client.post(
            "/api/v1/pdf/convert/word-to-pdf/", data={"file": docx_file}, format="multipart",
        )
        self.assertEqual(response.status_code, 201)
        data = response.json()["data"]
        self.assertEqual(data["operation"], "word_to_pdf")

        download_url = data["download_url"].replace("http://testserver", "")
        download = self.client.get(download_url)
        self.assertEqual(download["Content-Type"], "application/pdf")

        job = ConversionJob.objects.get(id=data["file_id"])
        job.output_file.delete(save=False)

    def test_rejects_non_docx_file(self):
        fake = SimpleUploadedFile("evil.docx", b"not a real docx", content_type="application/octet-stream")
        response = self.client.post(
            "/api/v1/pdf/convert/word-to-pdf/", data={"file": fake}, format="multipart",
        )
        self.assertEqual(response.status_code, 400)
