"""
PDF -> PPTX. One slide per page, reconstructed from real content (not a
flattened screenshot): text blocks become editable textboxes at their
original position, detected tables become real PowerPoint tables, and
images are placed at their actual on-page position/size. A page with no
extractable text (scanned/image-only) falls back to a full-slide image,
same reasoning as the Word/converter's scanned-page handling.

All slides share one presentation-wide canvas size (taken from the first
page). A later page with a different size is scaled per-axis to fit that
canvas - correct for the common case (uniform page size) and a documented
simplification for mixed-size documents.
"""

import io

import fitz
from pptx import Presentation
from pptx.util import Emu, Pt

BOLD_FLAG = 1 << 4
ITALIC_FLAG = 1 << 1
MIN_TEXT_LENGTH_FOR_REAL_CONVERSION = 5
POINTS_TO_EMU = 12700
BLANK_LAYOUT_INDEX = 6  # python-pptx's default template's "Blank" layout


def _pt_to_emu(value_pt):
    return Emu(int(round(value_pt * POINTS_TO_EMU)))


def _add_page_image_slide(prs, page, slide_width_emu, slide_height_emu):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    pix = page.get_pixmap(dpi=150)
    slide.shapes.add_picture(io.BytesIO(pix.tobytes("png")), 0, 0, width=slide_width_emu, height=slide_height_emu)
    return slide


def _add_textbox(slide, block, scale_x, scale_y):
    x0, y0, x1, y1 = block["bbox"]
    width = max(x1 - x0, 10)
    height = max(y1 - y0, 10)
    textbox = slide.shapes.add_textbox(
        _pt_to_emu(x0 * scale_x), _pt_to_emu(y0 * scale_y),
        _pt_to_emu(width * scale_x), _pt_to_emu(height * scale_y),
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    first_line = True
    for line in block.get("lines", []):
        paragraph = text_frame.paragraphs[0] if first_line else text_frame.add_paragraph()
        first_line = False
        for span in line.get("spans", []):
            text = span.get("text", "")
            if not text:
                continue
            run = paragraph.add_run()
            run.text = text
            run.font.bold = bool(span.get("flags", 0) & BOLD_FLAG)
            run.font.italic = bool(span.get("flags", 0) & ITALIC_FLAG)
            size = span.get("size")
            if size:
                run.font.size = Pt(round(size * min(scale_x, scale_y)))


def _add_table(slide, rows, bbox, scale_x, scale_y):
    if not rows or not rows[0]:
        return
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    x0, y0, x1, y1 = bbox
    width = max((x1 - x0) * scale_x, 10)
    height = max((y1 - y0) * scale_y, 10)
    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, _pt_to_emu(x0 * scale_x), _pt_to_emu(y0 * scale_y),
        _pt_to_emu(width), _pt_to_emu(height),
    )
    table = graphic_frame.table
    for row_index, row in enumerate(rows):
        for col_index in range(n_cols):
            value = row[col_index] if col_index < len(row) and row[col_index] is not None else ""
            table.cell(row_index, col_index).text = str(value)


def _rects_overlap_significantly(a, b, threshold=0.5):
    intersection = a & b
    if intersection.is_empty:
        return False
    a_area = a.get_area()
    if a_area == 0:
        return False
    return (intersection.get_area() / a_area) >= threshold


def convert(file_bytes):
    """Returns (pptx_bytes, metadata dict)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        prs = Presentation()
        page_count = len(doc)
        table_count = 0
        image_count = 0
        scanned_pages = []

        first_page = doc[0]
        prs.slide_width = _pt_to_emu(first_page.rect.width)
        prs.slide_height = _pt_to_emu(first_page.rect.height)

        for page_index, page in enumerate(doc):
            scale_x = prs.slide_width / _pt_to_emu(page.rect.width)
            scale_y = prs.slide_height / _pt_to_emu(page.rect.height)

            page_text = page.get_text().strip()
            if len(page_text) < MIN_TEXT_LENGTH_FOR_REAL_CONVERSION:
                scanned_pages.append(page_index + 1)
                _add_page_image_slide(prs, page, prs.slide_width, prs.slide_height)
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

            table_finder = page.find_tables()
            table_bboxes = [fitz.Rect(t.bbox) for t in table_finder.tables]

            text_dict = page.get_text("dict")
            blocks = sorted(
                (b for b in text_dict["blocks"] if b.get("type") == 0 and b.get("lines")),
                key=lambda b: (round(b["bbox"][1]), b["bbox"][0]),
            )
            for block in blocks:
                block_rect = fitz.Rect(block["bbox"])
                if any(_rects_overlap_significantly(block_rect, t) for t in table_bboxes):
                    continue
                _add_textbox(slide, block, scale_x, scale_y)

            for table in table_finder.tables:
                _add_table(slide, table.extract(), table.bbox, scale_x, scale_y)
                table_count += 1

            for image in page.get_images(full=True):
                xref = image[0]
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                try:
                    extracted = doc.extract_image(xref)
                except Exception:
                    continue
                rect = rects[0]
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(extracted["image"]),
                        _pt_to_emu(rect.x0 * scale_x), _pt_to_emu(rect.y0 * scale_y),
                        _pt_to_emu(rect.width * scale_x), _pt_to_emu(rect.height * scale_y),
                    )
                    image_count += 1
                except Exception:
                    continue

        buffer = io.BytesIO()
        prs.save(buffer)

        metadata = {
            "page_count": page_count,
            "table_count": table_count,
            "image_count": image_count,
            "scanned_pages": scanned_pages,
        }
        return buffer.getvalue(), metadata
    finally:
        doc.close()
